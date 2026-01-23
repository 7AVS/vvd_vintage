"""
Vintage Layer View - PowerPoint Generator
Generates both Executive Brief and Technical Onboarding presentations
Following the MVD → Visual Framework → PPTX pipeline
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================================================
# COLOR PALETTE (Minimalist Philosophy)
# ============================================================================
COLORS = {
    # Structural
    "WHITE": RGBColor(0xFF, 0xFF, 0xFF),
    "LIGHT_GRAY": RGBColor(0xF5, 0xF5, 0xF5),
    "BORDER_GRAY": RGBColor(0xCC, 0xCC, 0xCC),
    "DARK_BLUE": RGBColor(0x00, 0x31, 0x68),
    "MID_GRAY": RGBColor(0x66, 0x66, 0x66),

    # Accents (use sparingly)
    "OCEAN": RGBColor(0x00, 0x91, 0xDA),       # From Data (L1, L4)
    "WARM_YELLOW": RGBColor(0xFF, 0xC7, 0x2C), # Swap Points (L2, L3)
    "SUNBURST": RGBColor(0xFC, 0xA3, 0x11),    # Key Emphasis
    "TUNDRA": RGBColor(0x07, 0xAF, 0xBF),      # Future State
    "ERROR_RED": RGBColor(0xE5, 0x4B, 0x4B),   # Problem indicator
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Add a clean title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS["DARK_BLUE"]
    p.alignment = PP_ALIGN.CENTER

    # Line separator
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4), Inches(2), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["BORDER_GRAY"]
    line.line.fill.background()

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS["MID_GRAY"]
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title):
    """Add a slide with title, return the slide for further customization"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["DARK_BLUE"]

    return slide


def add_box(slide, left, top, width, height, text="", fill_color=None,
            border_color=None, text_color=None, font_size=12, bold=False,
            italic=False, align=PP_ALIGN.LEFT, monospace=False):
    """Add a box with optional fill, border, and text"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))

    # Fill
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    # Border
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()

    # Text
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = text_color or COLORS["DARK_BLUE"]
        p.alignment = align
        if monospace:
            p.font.name = "Consolas"

        # Vertical center
        tf.anchor = MSO_ANCHOR.MIDDLE

    return shape


def add_text(slide, left, top, width, height, text, font_size=12,
             color=None, bold=False, italic=False, align=PP_ALIGN.LEFT):
    """Add a text box"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color or COLORS["DARK_BLUE"]
    p.alignment = align
    return box


def add_arrow(slide, start_left, start_top, end_left, end_top, color=None):
    """Add a connector arrow"""
    # Using a simple line shape
    connector = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(start_left), Inches(start_top),
        Inches(0.3), Inches(end_top - start_top)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color or COLORS["BORDER_GRAY"]
    connector.line.fill.background()
    return connector


# ============================================================================
# EXECUTIVE BRIEF SLIDES
# ============================================================================

def create_executive_brief():
    """Generate the Executive Brief presentation (9 slides)"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "VINTAGE AUTOMATION",
                    "Building Measurement Infrastructure That Scales\n\nExecutive Brief")

    # Slide 2: The Problem
    slide = add_content_slide(prs, "The Problem We Solved")

    # Problem container
    add_box(slide, 0.5, 1.2, 9, 1.5,
            "TODAY: Every measurement request means...",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
            font_size=14, bold=True)

    # Three problem boxes
    problems = ["Manual Data\nExtraction", "Tribal Knowledge\n\"How to calculate\"", "Weeks of Work\nper Campaign"]
    for i, prob in enumerate(problems):
        add_box(slide, 1 + i*2.8, 2.9, 2.4, 1.2, prob,
                border_color=COLORS["BORDER_GRAY"], font_size=12, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 4.3, 9, 0.5, "Result: Inconsistent definitions across teams",
             font_size=14, italic=True, color=COLORS["MID_GRAY"])

    # Key callout
    add_box(slide, 1.5, 5.2, 7, 0.8, "WE BUILT AN ENGINE THAT CHANGES THIS",
            border_color=COLORS["SUNBURST"], font_size=18, bold=True, align=PP_ALIGN.CENTER)

    # Slide 3: 4-Layer Architecture
    slide = add_content_slide(prs, "The Vintage Automation Engine")
    add_text(slide, 0.5, 0.8, 9, 0.4, "SuperFact 4-Layer Framework", font_size=14, color=COLORS["MID_GRAY"])

    # Four layer boxes
    layers = [
        ("Layer 1\nWho is in test?", COLORS["OCEAN"]),
        ("Layer 2\nWhat metric?", COLORS["WARM_YELLOW"]),
        ("Layer 3\nHow to calc?", COLORS["WARM_YELLOW"]),
        ("Layer 4\nWhat did they do?", COLORS["OCEAN"])
    ]
    for i, (label, color) in enumerate(layers):
        add_box(slide, 0.6 + i*2.3, 1.4, 2.1, 1.1, label,
                border_color=color, font_size=11, align=PP_ALIGN.CENTER)

    # Arrow down
    add_arrow(slide, 4.8, 2.6, 4.8, 3.2, COLORS["BORDER_GRAY"])

    # Engine box
    add_box(slide, 2.5, 3.3, 5, 1, "VINTAGE ENGINE\n(Stable Core)",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["DARK_BLUE"],
            font_size=14, bold=True, align=PP_ALIGN.CENTER)

    # Arrow down to tracks
    add_arrow(slide, 4.8, 4.4, 4.8, 5.0, COLORS["BORDER_GRAY"])

    # Two tracks
    add_box(slide, 1.5, 5.1, 2.8, 1, "TRACK A\nTableau",
            border_color=COLORS["BORDER_GRAY"], font_size=12, align=PP_ALIGN.CENTER)
    add_text(slide, 4.5, 5.3, 1, 0.5, "PARALLEL", font_size=10, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 5.7, 5.1, 2.8, 1, "TRACK B\nSharePoint",
            border_color=COLORS["BORDER_GRAY"], font_size=12, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 6.3, 9, 0.5, "One engine. Two delivery tracks. Running in parallel.",
             font_size=14, bold=True, align=PP_ALIGN.CENTER)

    # Slide 4: Current State
    slide = add_content_slide(prs, "Current State: Proven and Ready")

    # Status boxes
    stats = [("6", "Campaigns\nMeasured", "VCN VDA VDT\nVUI VUT VAW"),
             ("4", "Metrics\nDefined", "card_acq\ncard_act\ncard_use\nwallet_prov"),
             ("✓", "Engine\nBuilt", "")]
    for i, (num, label, detail) in enumerate(stats):
        add_box(slide, 0.8 + i*3, 1.2, 2.6, 2.5, f"{num}\n{label}\n\n{detail}",
                border_color=COLORS["BORDER_GRAY"], font_size=12, align=PP_ALIGN.CENTER)

    # Track status
    add_box(slide, 0.8, 4.2, 4, 1, "Track A: Pending\nAlignment with CIDM",
            border_color=COLORS["BORDER_GRAY"], font_size=12, align=PP_ALIGN.CENTER)
    add_box(slide, 5.2, 4.2, 4, 1, "Track B: Ready NOW\nHTML on SharePoint",
            border_color=COLORS["TUNDRA"], font_size=12, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 5.5, 9, 0.5, "We're not waiting. Track B delivers value TODAY.",
             font_size=14, bold=True, align=PP_ALIGN.CENTER)

    # Slide 5: Virtuous Cycle
    slide = add_content_slide(prs, "The Virtuous Cycle: Why This Scales")

    # Cycle boxes
    cycle_labels = ["New\nCampaign", "Define\nMetric", "Add to\nLibrary", "Next\nFaster"]
    positions = [(1, 1.5), (3.5, 1.5), (6, 1.5), (3.5, 3.2)]
    for i, ((x, y), label) in enumerate(zip(positions, cycle_labels)):
        color = COLORS["SUNBURST"] if i == 2 else (COLORS["TUNDRA"] if i == 3 else COLORS["BORDER_GRAY"])
        add_box(slide, x, y, 2, 1, label, border_color=color, font_size=12, align=PP_ALIGN.CENTER)

    # Growth section
    add_box(slide, 0.5, 4.4, 9, 2.2,
            "GROWTH: 6 campaigns, only 4 unique metrics\n\n"
            "VCN → defines card_acq\n"
            "VDA → reuses card_acq ✓\n"
            "VDT → adds card_act\n"
            "VAW → reuses wallet_prov ✓\n\n"
            "Campaign 7, 8, 9... become trivial if metrics exist",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
            font_size=11, align=PP_ALIGN.LEFT)

    # Slide 6: Swap Points
    slide = add_content_slide(prs, "Modular Architecture: Built for the Future")
    add_text(slide, 0.5, 0.8, 9, 0.4, "59 hardcoded items → Dynamic when infrastructure ready",
             font_size=12, color=COLORS["MID_GRAY"])

    # Layer 2 swap
    add_box(slide, 0.5, 1.4, 9, 2, "LAYER 2: Campaign Metadata",
            border_color=COLORS["WARM_YELLOW"], font_size=11)
    add_box(slide, 0.8, 1.9, 3.5, 1.2, "TODAY\nPython dict\n6 campaigns hardcoded",
            border_color=COLORS["BORDER_GRAY"], font_size=10, align=PP_ALIGN.CENTER)
    add_text(slide, 4.4, 2.3, 1, 0.5, "→ SWAP →", font_size=10, bold=True, color=COLORS["SUNBURST"])
    add_box(slide, 5.5, 1.9, 3.5, 1.2, "FUTURE\nQuery\nMnemonic Mapping v2",
            border_color=COLORS["TUNDRA"], font_size=10, align=PP_ALIGN.CENTER)

    # Layer 3 swap
    add_box(slide, 0.5, 3.7, 9, 2, "LAYER 3: Success Definitions",
            border_color=COLORS["WARM_YELLOW"], font_size=11)
    add_box(slide, 0.8, 4.2, 3.5, 1.2, "TODAY\nFilters in code\n(24 items)",
            border_color=COLORS["BORDER_GRAY"], font_size=10, align=PP_ALIGN.CENTER)
    add_text(slide, 4.4, 4.6, 1, 0.5, "→ SWAP →", font_size=10, bold=True, color=COLORS["SUNBURST"])
    add_box(slide, 5.5, 4.2, 3.5, 1.2, "FUTURE\nSuccess Library\n(GitHub)",
            border_color=COLORS["TUNDRA"], font_size=10, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 6, 9, 0.5, "Engine core doesn't change - only the inputs swap",
             font_size=12, bold=True, align=PP_ALIGN.CENTER)

    # Slide 7: Two Tracks
    slide = add_content_slide(prs, "Two Tracks: Speed AND Governance")
    add_text(slide, 0.5, 0.8, 9, 0.4, "One engine feeds both. We're not choosing - we get BOTH.",
             font_size=12, color=COLORS["MID_GRAY"])

    # Track A
    add_box(slide, 0.5, 1.4, 4.3, 4.3,
            "TRACK A\nOfficial\n\n"
            "Platform:\nTableau via CIDM\n\n"
            "Status:\nPending alignment\n\n"
            "Strength:\nGoverned, trusted\n\n"
            "Refresh:\nAutomated (when ready)",
            border_color=COLORS["BORDER_GRAY"], font_size=11, align=PP_ALIGN.LEFT)

    # Track B
    add_box(slide, 5.2, 1.4, 4.3, 4.3,
            "TRACK B\nIn-House\n\n"
            "Platform:\nHTML/Plotly + SharePoint\n\n"
            "Status:\nREADY NOW\n\n"
            "Strength:\nFast, controlled\n\n"
            "Refresh:\nManual re-run",
            border_color=COLORS["BORDER_GRAY"], font_size=11, align=PP_ALIGN.LEFT)

    add_text(slide, 3.5, 5.9, 3, 0.5, "◀─── PARALLEL ───▶",
             font_size=12, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 6.5, 9, 0.5,
             "Track B proves value NOW. Track A becomes official when ready.",
             font_size=12, bold=True, align=PP_ALIGN.CENTER)

    # Slide 8: What We Need
    slide = add_content_slide(prs, "What We Need: Support and Buy-In")

    asks = [
        ("LEADERSHIP", "• Awareness\n• Support\n• Champion this\n  as standard"),
        ("CIDM /\nINFRASTRUCTURE", "• Alignment\n• Prioritize\n  integration\n• Collaborate"),
        ("OTHER\nTEAMS", "• Adoption\n• Contribute\n  definitions\n• Use Success\n  Library")
    ]
    for i, (header, content) in enumerate(asks):
        add_box(slide, 0.6 + i*3.1, 1.2, 2.9, 0.9, header,
                fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
                font_size=12, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, 0.6 + i*3.1, 2.1, 2.9, 3.5, content,
                border_color=COLORS["BORDER_GRAY"], font_size=11, align=PP_ALIGN.LEFT)

    # Slide 9: Vision
    slide = add_content_slide(prs, "The Vision: Self-Service Measurement")

    add_box(slide, 0.5, 1.2, 9, 3.5,
            "FUTURE STATE\n\n"
            "Marketing wants to measure a new campaign:\n\n"
            "1. Look up metric in Success Library      ✓ Already exists\n"
            "2. Add row to Mnemonic Mapping v2         ✓ Self-service\n"
            "3. Run Vintage Engine                     ✓ No code change\n"
            "4. View results in dashboard              ✓ Same day\n\n"
            "────────────────────────────────────────────────\n"
            "Total effort: MINUTES, not weeks\n"
            "Engineering involvement: ZERO (for existing metrics)",
            border_color=COLORS["SUNBURST"], font_size=12, align=PP_ALIGN.LEFT)

    add_box(slide, 0.5, 5, 9, 1.5,
            "\"Every campaign onboarded enriches our metadata ecosystem.\n"
            "We're not just measuring - we're building the source of truth.\"",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
            font_size=14, italic=True, align=PP_ALIGN.CENTER)

    return prs


# ============================================================================
# TECHNICAL ONBOARDING SLIDES
# ============================================================================

def create_technical_onboarding():
    """Generate the Technical Onboarding presentation (10 slides)"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "VINTAGE AUTOMATION",
                    "Technical Onboarding Guide\n\nJoin Us in Building the Measurement Standard")

    # Slide 2: Why This Matters to YOU
    slide = add_content_slide(prs, "Why This Matters to YOU")
    add_text(slide, 0.5, 0.9, 9, 0.4, "You've probably done this before:",
             font_size=14, color=COLORS["MID_GRAY"])

    add_box(slide, 0.5, 1.4, 9, 3.2,
            "Someone asks: \"How many cards acquired for campaign X?\"\n\n"
            "        You                            Someone else\n"
            "    ┌─────────────┐                 ┌─────────────┐\n"
            "    │ Write SQL   │                 │ Write SQL   │\n"
            "    │ Find tables │                 │ Find tables │\n"
            "    │ Add filters │                 │ Add filters │\n"
            "    └──────┬──────┘                 └──────┬──────┘\n"
            "           ▼                               ▼\n"
            "      Answer A            ≠           Answer B\n"
            "   (Different filters)              (Different filters)\n\n"
            "Same question. Different answers.",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
            font_size=11, monospace=True)

    add_box(slide, 1.5, 4.9, 7, 0.8, "We're fixing this. And we need YOUR help.",
            border_color=COLORS["SUNBURST"], font_size=16, bold=True, align=PP_ALIGN.CENTER)

    # Slide 3: 4-Layer Architecture (Detailed)
    slide = add_content_slide(prs, "The Architecture: SuperFact 4-Layer Framework")
    add_text(slide, 0.5, 0.8, 9, 0.4, "vintage_all_in_one.py", font_size=12,
             color=COLORS["MID_GRAY"], italic=True)

    # Layer boxes
    layers = [
        ("LAYER 1: EXPERIMENT METADATA", "\"Who is in the test?\"",
         "Code: load_tactic() | Config: YEARS, TEST_GROUP_CODE",
         "Swap to: → Experiment Metadata table", COLORS["OCEAN"], 1.2),
        ("LAYER 2: CAMPAIGN METADATA", "\"What metric to measure?\"",
         "Code: get_campaign_config() | Config: CAMPAIGN_METADATA",
         "Swap to: → Mnemonic Mapping v2 query", COLORS["WARM_YELLOW"], 2.5),
        ("LAYER 3: SUCCESS DEFINITIONS", "\"How to calculate?\"",
         "Code: get_success_definition() | Config: SUCCESS_DEFS",
         "Swap to: → Success Library (GitHub or curated data)", COLORS["WARM_YELLOW"], 3.8),
        ("LAYER 4: CLIENT JOURNEY", "\"What did they actually do?\"",
         "Code: load_fulfillment(), load_email_engagement(),...",
         "Swap to: → Unified Client Journey semantic layer", COLORS["OCEAN"], 5.1)
    ]

    for header, question, code, swap, color, top in layers:
        add_box(slide, 0.5, top, 9, 1.2, f"{header}\n{question}\n{code}\n{swap}",
                border_color=color, font_size=10, monospace=False)

    # Slide 4: Engine + Two Tracks
    slide = add_content_slide(prs, "Engine + Delivery Tracks")

    add_text(slide, 4, 1, 2, 0.4, "4 LAYERS", font_size=12, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.8, 1.3, 4.8, 1.8, COLORS["BORDER_GRAY"])

    add_box(slide, 1.5, 1.9, 7, 1.8,
            "VINTAGE ENGINE (Stable - does NOT change)\n\n"
            "detect_success() → build_vintage_data()\n"
            "→ calculate_ci() → prepare_vintage_table()\n"
            "→ generate_summary() → plot_vintage()\n\n"
            "This block doesn't care WHERE data comes from.",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["DARK_BLUE"],
            font_size=11, monospace=True)

    add_arrow(slide, 4.8, 3.8, 4.8, 4.3, COLORS["BORDER_GRAY"])
    add_text(slide, 4, 4.1, 2, 0.4, "CSV / HDFS", font_size=10, align=PP_ALIGN.CENTER)

    add_box(slide, 1, 4.6, 3.5, 1.3, "TRACK A\nTableau / CIDM\n(pending)",
            border_color=COLORS["BORDER_GRAY"], font_size=12, align=PP_ALIGN.CENTER)
    add_text(slide, 4.6, 5, 1, 0.5, "PARALLEL", font_size=10, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 5.5, 4.6, 3.5, 1.3, "TRACK B\nSharePoint / HTML/Plotly\n(ready)",
            border_color=COLORS["BORDER_GRAY"], font_size=12, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 6.2, 9, 0.5, "Same engine feeds BOTH. No duplicate work.",
             font_size=14, bold=True, align=PP_ALIGN.CENTER)

    # Slide 5: Swap Point - Layer 2
    slide = add_content_slide(prs, "Swap Point: Layer 2 - Campaign Metadata")

    # Today box
    add_box(slide, 0.5, 1.1, 9, 2,
            "TODAY (hardcoded)\n\n"
            "CAMPAIGN_METADATA = {\n"
            "    \"VCN\": {\n"
            "        \"campaign_name\": \"VVD Contextual...\",\n"
            "        \"success_type\": \"ACQUISITION\",\n"
            "        \"primary_metric\": \"card_acquisition\",\n"
            "    },\n"
            "    # ... 6 campaigns hardcoded\n"
            "}",
            border_color=COLORS["WARM_YELLOW"], font_size=10, monospace=True)

    add_text(slide, 4.5, 3.2, 1, 0.5, "SWAP ▼", font_size=12, bold=True,
             color=COLORS["SUNBURST"], align=PP_ALIGN.CENTER)

    # Future box
    add_box(slide, 0.5, 3.7, 9, 1.5,
            "FUTURE (dynamic)\n\n"
            "config = spark.sql(\"\"\"\n"
            "    SELECT primary_metric, secondary_metric\n"
            "    FROM mnemonic_mapping_v2 WHERE mne = 'VCN'\n"
            "\"\"\")",
            border_color=COLORS["TUNDRA"], font_size=10, monospace=True)

    # Your opportunity
    add_box(slide, 0.5, 5.5, 9, 1,
            "YOUR OPPORTUNITY: Enhancing MM v2 with metric fields?\n"
            "The engine will automatically use your work.",
            fill_color=COLORS["SUNBURST"], border_color=COLORS["SUNBURST"],
            text_color=COLORS["WHITE"], font_size=12, bold=True, align=PP_ALIGN.CENTER)

    # Slide 6: Swap Point - Layer 3
    slide = add_content_slide(prs, "Swap Point: Layer 3 - Success Definitions (THE BIG ONE)")

    # Today box
    add_box(slide, 0.5, 1, 9, 1.8,
            "TODAY\n"
            "SUCCESS_DEFINITIONS = {\n"
            "    \"card_acquisition\": {\n"
            "        \"table_path\": \"/prod/.../VISA_DR_CRD/\",\n"
            "        \"filters\": {\"STS_CD\": [\"06\",\"08\"], \"SRVC_ID\": 36},\n"
            "    }, # ... 4 metrics × 6 fields = 24 items\n"
            "}",
            border_color=COLORS["WARM_YELLOW"], font_size=10, monospace=True)

    add_text(slide, 4.5, 2.9, 1, 0.5, "SWAP ▼", font_size=12, bold=True,
             color=COLORS["SUNBURST"], align=PP_ALIGN.CENTER)

    # Future options
    add_box(slide, 0.5, 3.4, 4.3, 1.6,
            "FUTURE: Option A - GitHub %Run\n\n"
            "%run /success_lib/card_acquisition.py\n"
            "success_df = get_card_acquisition(\n"
            "    spark, client_list)",
            border_color=COLORS["TUNDRA"], font_size=10, monospace=True)

    add_box(slide, 5.2, 3.4, 4.3, 1.6,
            "FUTURE: Option B - Curated Data\n\n"
            "spark.read.parquet(\n"
            "    \"/semantic/card_acquisition\"\n"
            ")",
            border_color=COLORS["TUNDRA"], font_size=10, monospace=True)

    # Your opportunity
    add_box(slide, 0.5, 5.3, 9, 1.2,
            "YOUR OPPORTUNITY: Have metric definitions? Curated data?\n"
            "Share them! They become the standard.",
            fill_color=COLORS["SUNBURST"], border_color=COLORS["SUNBURST"],
            text_color=COLORS["WHITE"], font_size=12, bold=True, align=PP_ALIGN.CENTER)

    # Slide 7: Virtuous Cycle
    slide = add_content_slide(prs, "The Virtuous Cycle: Why Your Contribution Matters")

    add_box(slide, 0.5, 1.1, 9, 3.5,
            "You define \"card_acquisition\" once\n"
            "              │\n"
            "              ▼\n"
            "It goes in the Success Library\n"
            "              │\n"
            "              ▼\n"
            "Campaign 1 uses it ✓\n"
            "Campaign 2 uses it ✓  (no new work!)\n"
            "Campaign 5 uses it ✓  (still no new work!)\n"
            "              │\n"
            "              ▼\n"
            "Another team needs \"card_acquisition\"?\n"
            "They use YOUR definition.\n"
            "Consistency. Trust. YOUR name on it.",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
            font_size=11, monospace=True)

    add_box(slide, 0.5, 4.9, 9, 1.2,
            "RESULT: Define once, reuse forever\n"
            "YOUR CONTRIBUTION: Multiplied across the organization",
            border_color=COLORS["SUNBURST"], font_size=14, bold=True, align=PP_ALIGN.CENTER)

    # Slide 8: 59 Swap Points
    slide = add_content_slide(prs, "What's Hardcoded Today: 59 Items Ready to Swap")

    # Table header
    add_box(slide, 0.5, 1.1, 2, 0.5, "Layer", fill_color=COLORS["LIGHT_GRAY"],
            border_color=COLORS["BORDER_GRAY"], font_size=11, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 2.5, 1.1, 3.5, 0.5, "Hardcoded Items", fill_color=COLORS["LIGHT_GRAY"],
            border_color=COLORS["BORDER_GRAY"], font_size=11, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 6, 1.1, 1, 0.5, "Count", fill_color=COLORS["LIGHT_GRAY"],
            border_color=COLORS["BORDER_GRAY"], font_size=11, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 7, 1.1, 2.5, 0.5, "Swaps To", fill_color=COLORS["LIGHT_GRAY"],
            border_color=COLORS["BORDER_GRAY"], font_size=11, bold=True, align=PP_ALIGN.CENTER)

    # Table rows
    rows = [
        ("Layer 1", "Years, test group, path", "3", "Experiment Metadata", COLORS["OCEAN"]),
        ("Layer 2", "Campaign config (6×4)", "24", "Mnemonic Mapping v2", COLORS["WARM_YELLOW"]),
        ("Layer 3", "Success defs (4×6)", "24", "Success Library", COLORS["WARM_YELLOW"]),
        ("Layer 4", "Paths, EDW queries", "8", "Semantic Layers", COLORS["OCEAN"]),
        ("TOTAL", "", "59", "", COLORS["DARK_BLUE"])
    ]

    for i, (layer, items, count, swaps, color) in enumerate(rows):
        top = 1.6 + i * 0.7
        bold = i == 4
        add_box(slide, 0.5, top, 2, 0.7, layer, border_color=color,
                font_size=11, bold=bold, align=PP_ALIGN.CENTER)
        add_box(slide, 2.5, top, 3.5, 0.7, items, border_color=COLORS["BORDER_GRAY"],
                font_size=10, align=PP_ALIGN.CENTER)
        add_box(slide, 6, top, 1, 0.7, count, border_color=COLORS["BORDER_GRAY"],
                font_size=11, bold=bold, align=PP_ALIGN.CENTER)
        add_box(slide, 7, top, 2.5, 0.7, swaps, border_color=COLORS["BORDER_GRAY"],
                font_size=10, align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 5.5, 9, 0.5, "59 items ready to swap when YOUR work is ready.",
             font_size=14, bold=True, color=COLORS["SUNBURST"], align=PP_ALIGN.CENTER)

    # Slide 9: How to Contribute
    slide = add_content_slide(prs, "How to Contribute")

    contributions = [
        ("Metric definition", "Document filters + tables, share with us\n→ Goes in Success Library"),
        ("Curated data set", "Share path + schema\n→ Engine points to it"),
        ("Semantic layer", "Tell us when ready\n→ We flip the switch"),
        ("\"The right way\" to calculate", "Tell us! Document it\n→ Tribal knowledge becomes standard")
    ]

    for i, (have, do) in enumerate(contributions):
        top = 1.1 + i * 1.3
        add_box(slide, 0.5, top, 3.5, 1.1, f"IF YOU HAVE...\n\n{have}",
                fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
                font_size=10, bold=False, align=PP_ALIGN.LEFT)
        add_box(slide, 4.2, top, 5.3, 1.1, f"DO THIS...\n\n{do}",
                border_color=COLORS["BORDER_GRAY"], font_size=10, align=PP_ALIGN.LEFT)

    # Slide 10: The Ask
    slide = add_content_slide(prs, "The Ask: Join Us")

    add_box(slide, 0.5, 1.1, 9, 2.5, "WE NEED YOU TO:",
            fill_color=COLORS["LIGHT_GRAY"], border_color=COLORS["BORDER_GRAY"],
            font_size=12, bold=True)

    asks = [("SHARE", "Your metric\ndefinitions"),
            ("COLLABORATE", "Fill gaps in\nSuccess Library"),
            ("ADOPT", "Use standard\ninstead of own"),
            ("CONTRIBUTE", "Your semantic\nlayers plug in")]
    for i, (header, content) in enumerate(asks):
        add_box(slide, 0.8 + i*2.2, 1.8, 2, 1.5, f"{header}\n\n{content}",
                border_color=COLORS["BORDER_GRAY"], font_size=10, align=PP_ALIGN.CENTER)

    add_box(slide, 0.5, 3.9, 9, 2,
            "WHAT YOU GET:\n\n"
            "• Your definitions become THE standard\n"
            "• Your work gets reused across campaigns\n"
            "• Less \"can you pull this data?\" requests\n"
            "• Consistency across the entire team",
            border_color=COLORS["TUNDRA"], font_size=12, align=PP_ALIGN.LEFT)

    add_text(slide, 0.5, 6.2, 9, 0.7,
             "\"We're not just measuring campaigns. We're building the\n"
             "metadata ecosystem that makes measurement easy for everyone.\"",
             font_size=12, italic=True, align=PP_ALIGN.CENTER)

    return prs


# ============================================================================
# MAIN
# ============================================================================

def main():
    """Generate both presentations"""
    output_dir = os.path.dirname(os.path.abspath(__file__))

    # Executive Brief
    print("Generating Executive Brief...")
    exec_prs = create_executive_brief()
    exec_path = os.path.join(output_dir, "EXECUTIVE_BRIEF.pptx")
    exec_prs.save(exec_path)
    print(f"  Saved: {exec_path}")

    # Technical Onboarding
    print("Generating Technical Onboarding...")
    tech_prs = create_technical_onboarding()
    tech_path = os.path.join(output_dir, "TECHNICAL_ONBOARDING.pptx")
    tech_prs.save(tech_path)
    print(f"  Saved: {tech_path}")

    print("\nDone! Both presentations generated.")


if __name__ == "__main__":
    main()
