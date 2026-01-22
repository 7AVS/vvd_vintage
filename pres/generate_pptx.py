"""
Generate Vintage Automation PowerPoint from MVD content
Creates native editable shapes - no images required
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# RBC Brand Colors
DARK_BLUE = RGBColor(0x00, 0x31, 0x68)
BRIGHT_BLUE = RGBColor(0x00, 0x51, 0xA5)
OCEAN = RGBColor(0x00, 0x91, 0xDA)
WARM_YELLOW = RGBColor(0xFF, 0xC7, 0x2C)
TUNDRA = RGBColor(0x07, 0xAF, 0xBF)
SUNBURST = RGBColor(0xFC, 0xA3, 0x11)
GRAY = RGBColor(0x9E, 0xA2, 0xA2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE7, 0xEE, 0xF1)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title box
    title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.33), Inches(1.5))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = DARK_BLUE
    title_box.line.fill.background()

    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_box(slide, left, top, width, height, fill_color, text="", font_size=11, font_color=WHITE, bold=False, border_color=None):
    """Add a colored box with text"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color

    return shape

def add_text(slide, left, top, width, height, text, font_size=11, font_color=DARK_BLUE, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.alignment = align
    return textbox

def create_page1_big_picture(prs):
    """Page 1: Big Picture - SuperFact Implementation"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text(slide, 0.5, 0.2, 12.33, 0.5, "Vintage Automation Engine - SuperFact Implementation",
             font_size=28, bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # SuperFact Header
    add_box(slide, 0.5, 0.8, 12.33, 0.4, DARK_BLUE, "SUPERFACT 4-LAYER FRAMEWORK", 14, WHITE, True)

    # Four Layers
    layer_width = 2.9
    layer_start = 0.6
    layer_gap = 3.1

    layers = [
        ("Layer 1", "Experiment Metadata", '"Who\'s in test?"', OCEAN),
        ("Layer 2", "Campaign Metadata", '"What to measure?"', OCEAN),
        ("Layer 3", "Success Library", '"How to calculate?"', OCEAN),
        ("Layer 4", "Client Journey", '"What did they do?"', OCEAN),
    ]

    for i, (title, subtitle, question, color) in enumerate(layers):
        left = layer_start + (i * layer_gap)
        box = add_box(slide, left, 1.3, layer_width, 0.9, color, f"{title}\n{subtitle}\n\n{question}", 10, WHITE, False, BRIGHT_BLUE)

    # Arrow down
    add_text(slide, 6.2, 2.25, 1, 0.3, "▼", 24, BRIGHT_BLUE, True, PP_ALIGN.CENTER)

    # Engine Container
    add_box(slide, 0.5, 2.6, 12.33, 2.8, WHITE, "", border_color=BRIGHT_BLUE)
    add_box(slide, 0.5, 2.6, 12.33, 0.35, BRIGHT_BLUE, "VINTAGE AUTOMATION ENGINE (vintage_all_in_one.py)", 12, WHITE, True)

    # Input boxes inside engine
    add_text(slide, 0.7, 3.05, 3, 0.25, "INPUTS FROM LAYERS:", 10, DARK_BLUE, True)

    inputs = [
        ("Layer 1\ntactic_evnt_hist\n\nFROM DATA", OCEAN, WHITE),
        ("Layer 2\nCampaign Config\n\nSWAP POINT", WARM_YELLOW, DARK_BLUE),
        ("Layer 3\nSuccess Definitions\n\nSWAP POINT", WARM_YELLOW, DARK_BLUE),
        ("Layer 4\nVISA_DR_CRD\nEmail Data\nFROM DATA", OCEAN, WHITE),
    ]

    for i, (text, bg, fg) in enumerate(inputs):
        add_box(slide, 0.7 + (i * 2.4), 3.35, 2.2, 0.85, bg, text, 9, fg)

    # Legend
    add_box(slide, 10.3, 3.35, 2.4, 0.85, WHITE, "", border_color=GRAY)
    add_text(slide, 10.4, 3.4, 2.2, 0.2, "LEGEND", 9, DARK_BLUE, True, PP_ALIGN.CENTER)
    add_box(slide, 10.5, 3.65, 0.25, 0.15, OCEAN)
    add_text(slide, 10.85, 3.62, 1.8, 0.2, "From Data (working)", 8, DARK_BLUE)
    add_box(slide, 10.5, 3.85, 0.25, 0.15, WARM_YELLOW)
    add_text(slide, 10.85, 3.82, 1.8, 0.2, "Swap Point", 8, DARK_BLUE)

    # Engine core
    add_text(slide, 0.7, 4.3, 6, 0.25, "ENGINE CORE (stable - doesn't change when sources change):", 10, DARK_BLUE, True)
    add_box(slide, 0.7, 4.6, 8, 0.35, BRIGHT_BLUE, "detect_success() → build_vintage_data() → calculate_ci() → generate_summary()", 10, WHITE, False)

    # Pilots and outputs
    add_text(slide, 0.7, 5.05, 6, 0.2, "PILOTS: VCN, VDA, VDT, VUI, VUT, VAW (6 campaigns)", 10, DARK_BLUE)
    add_text(slide, 0.7, 5.25, 8, 0.2, "OUTPUTS: Vintage curves (Test vs Control) | Lift + CI | Engagement funnel", 10, DARK_BLUE)

    # Arrow down
    add_text(slide, 6.2, 5.5, 1, 0.3, "▼", 24, BRIGHT_BLUE, True, PP_ALIGN.CENTER)

    # CSV Output
    add_box(slide, 5, 5.85, 3.33, 0.4, BRIGHT_BLUE, "CSV / HDFS OUTPUT", 12, WHITE, True)

    # Two tracks
    add_box(slide, 1.5, 6.5, 3.5, 1.1, GRAY, "TRACK A - Official\n\nWho: CIDM team\nTool: Tableau\nStatus: Pending", 10, WHITE)
    add_box(slide, 8.3, 6.5, 3.5, 1.1, OCEAN, "TRACK B - In-House\n\nWho: Us\nTool: HTML/Plotly\nStatus: Ready", 10, WHITE)

    return slide

def create_page2_virtuous_cycle(prs):
    """Page 2: Virtuous Cycle - Success Library Growth"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text(slide, 0.3, 0.15, 12.7, 0.4, "The Virtuous Cycle - Building the Source of Truth",
             font_size=26, bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.3, 0.5, 12.7, 0.3,
             "As we add campaigns, we document success metrics → Library grows → Governance improves → Next campaign is faster",
             font_size=12, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Left side: The Cycle
    add_box(slide, 0.2, 0.85, 4.2, 4.3, WHITE, "", border_color=BRIGHT_BLUE)
    add_box(slide, 0.2, 0.85, 4.2, 0.3, BRIGHT_BLUE, "THE CYCLE", 12, WHITE, True)

    steps = [
        ("1. NEW CAMPAIGN ONBOARDED", BRIGHT_BLUE, WHITE),
        ("2. UNDERSTAND SUCCESS METRIC", OCEAN, WHITE),
        ("3. DOCUMENT IN SUCCESS LIBRARY", WARM_YELLOW, DARK_BLUE),
        ("4. SUCCESS LIBRARY GROWS", TUNDRA, WHITE),
        ("5. GOVERNANCE IMPROVES", TUNDRA, WHITE),
        ("6. NEXT CAMPAIGN IS FASTER", SUNBURST, DARK_BLUE),
    ]

    for i, (text, bg, fg) in enumerate(steps):
        add_box(slide, 0.5, 1.25 + (i * 0.55), 3.6, 0.45, bg, text, 10, fg, True)
        if i < 5:
            add_text(slide, 2.1, 1.7 + (i * 0.55), 0.5, 0.2, "↓", 14, BRIGHT_BLUE, True, PP_ALIGN.CENTER)

    # Right side: Success Library Growth
    add_box(slide, 4.6, 0.85, 8.5, 4.3, WHITE, "", border_color=BRIGHT_BLUE)
    add_box(slide, 4.6, 0.85, 8.5, 0.3, BRIGHT_BLUE, "SUCCESS LIBRARY GROWTH OVER 6 CAMPAIGNS", 12, WHITE, True)

    campaigns = [
        ("VCN", "Card Acquisition", "DEFINE", ["card_acquisition"], [True]),
        ("VDA", "Card Acquisition", "REUSE", ["card_acquisition"], [False]),
        ("VDT", "Card Activation", "DEFINE", ["card_acquisition", "card_activation"], [False, True]),
        ("VUI", "Card Usage", "DEFINE", ["card_acquisition", "card_activation", "card_usage"], [False, False, True]),
        ("VUT", "Wallet Provisioning", "DEFINE", ["card_acquisition", "card_activation", "card_usage", "wallet_provision"], [False, False, False, True]),
        ("VAW", "Wallet Provisioning", "REUSE", ["card_acquisition", "card_activation", "card_usage", "wallet_provision"], [False, False, False, False]),
    ]

    y_start = 1.2
    for i, (mne, need, action, metrics, is_new) in enumerate(campaigns):
        y = y_start + (i * 0.52)
        # Background
        add_box(slide, 4.75, y, 8.2, 0.48, LIGHT_GRAY, "", border_color=BRIGHT_BLUE)
        # Campaign info
        add_text(slide, 4.85, y + 0.02, 4, 0.22, f"{mne}: {need} → {action}", 9, DARK_BLUE, True)
        # Metrics
        x_metric = 4.85
        for j, (metric, new) in enumerate(zip(metrics, is_new)):
            color = TUNDRA if new else OCEAN
            add_box(slide, x_metric + (j * 1.1), y + 0.24, 1.0, 0.2, color, metric[:12], 7, WHITE)
        # NEW/REUSED tag
        if action == "DEFINE":
            add_box(slide, 12.3, y + 0.12, 0.5, 0.18, SUNBURST, "NEW", 7, DARK_BLUE, True)
        else:
            add_box(slide, 12.1, y + 0.12, 0.7, 0.18, OCEAN, "REUSED", 7, WHITE, True)

    # Result box
    add_box(slide, 4.6, 5.3, 8.5, 0.6, SUNBURST,
            "RESULT: 6 campaigns onboarded, only 4 unique metrics defined\nFUTURE: Campaign 7, 8, 9... likely REUSE existing metrics",
            11, DARK_BLUE, True)

    # Key message box
    add_box(slide, 0.2, 6.1, 12.9, 0.8, DARK_BLUE,
            "KEY MESSAGE\n\n\"Every campaign onboarded enriches our metadata ecosystem.\nWe're not just measuring campaigns - we're BUILDING THE SOURCE OF TRUTH.\"",
            12, WHITE, True)

    return slide

def create_page3_swap_points(prs):
    """Page 3: Swap Points - What Swaps to What"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text(slide, 0.3, 0.1, 12.7, 0.35, "Modular Swap Architecture - What Swaps To What",
             font_size=24, bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.3, 0.4, 12.7, 0.25,
             "Each hardcoded component has a SPECIFIC future replacement. When infrastructure is ready, we swap - no engine rewrite.",
             font_size=11, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Layer 2 Section
    add_box(slide, 0.2, 0.7, 12.9, 1.8, WHITE, "", border_color=WARM_YELLOW)
    add_box(slide, 0.2, 0.7, 12.9, 0.28, WARM_YELLOW, "LAYER 2: Campaign Metadata — \"What metric should we measure?\"", 11, DARK_BLUE, True)

    # Today box
    add_box(slide, 0.3, 1.05, 5.8, 0.22, DARK_BLUE, "TODAY: HARDCODED", 10, WHITE, True)
    add_box(slide, 0.3, 1.3, 5.8, 1.1, WARM_YELLOW,
            "CAMPAIGN_METADATA = {\n  \"VCN\": {\"primary_metric\": \"card_acquisition\"},\n  \"VDA\": {...}, \"VDT\": {...},\n  \"VUI\": {...}, \"VUT\": {...}, \"VAW\": {...}\n}",
            9, DARK_BLUE)

    # Arrow
    add_text(slide, 6.2, 1.7, 0.8, 0.3, "SWAP →", 12, SUNBURST, True, PP_ALIGN.CENTER)

    # Future box
    add_box(slide, 7.1, 1.05, 5.9, 0.22, DARK_BLUE, "FUTURE: QUERY MNEMONIC MAPPING v2", 10, WHITE, True)
    add_box(slide, 7.1, 1.3, 5.9, 1.1, TUNDRA,
            "SELECT primary_metric, secondary_metric\nFROM CIDM_MNEMONIC_ATTRS\nWHERE mne = 'VCN'\n\nWhen MM v2 has metric fields → this dict becomes a query",
            9, WHITE)

    # Layer 3 Section
    add_box(slide, 0.2, 2.6, 12.9, 2.2, WHITE, "", border_color=WARM_YELLOW)
    add_box(slide, 0.2, 2.6, 12.9, 0.28, WARM_YELLOW, "LAYER 3: Success Definitions — \"HOW do we calculate this metric?\"", 11, DARK_BLUE, True)

    # Today box
    add_box(slide, 0.3, 2.95, 5.8, 0.22, DARK_BLUE, "TODAY: HARDCODED", 10, WHITE, True)
    add_box(slide, 0.3, 3.2, 5.8, 1.5, WARM_YELLOW,
            "SUCCESS_DEFINITIONS = {\n  \"card_acquisition\": {\n    \"table\": \"VISA_DR_CRD\",\n    \"filters\": {\"STS_CD\": [\"06\",\"08\"]},\n    \"date_field\": \"ISS_DT\"\n  },\n  \"card_activation\": {...},\n  \"card_usage\": {...}, ...\n}",
            8, DARK_BLUE)

    # Arrow
    add_text(slide, 6.2, 3.8, 0.8, 0.3, "SWAP →", 12, SUNBURST, True, PP_ALIGN.CENTER)

    # Future boxes
    add_box(slide, 7.1, 2.95, 5.9, 0.22, DARK_BLUE, "FUTURE: SUCCESS LIBRARY", 10, WHITE, True)
    add_box(slide, 7.1, 3.2, 2.85, 0.9, TUNDRA,
            "Option A: GitHub %Run\n\n%run /success_library/\n      card_acquisition.py",
            8, WHITE)
    add_box(slide, 10.05, 3.2, 2.95, 0.9, TUNDRA,
            "Option B: Curated Dataset\n\nSELECT * FROM\nsemantic.success_card_acq",
            8, WHITE)
    add_box(slide, 7.1, 4.15, 5.9, 0.55, LIGHT_GRAY,
            "4 metrics × 6 fields = 24 items today → GitHub repo OR semantic layer when ready",
            9, DARK_BLUE, border_color=BRIGHT_BLUE)

    # Summary Table
    add_box(slide, 0.2, 4.95, 12.9, 1.95, WHITE, "", border_color=BRIGHT_BLUE)
    add_box(slide, 0.2, 4.95, 12.9, 0.28, BRIGHT_BLUE, "SWAP POINTS SUMMARY", 11, WHITE, True)

    # Table headers
    headers = [("LAYER", 1.2), ("WHAT'S HARDCODED", 3.5), ("SWAPS TO", 3.5), ("TRIGGER", 3.8)]
    x = 0.35
    for header, width in headers:
        add_box(slide, x, 5.28, width, 0.25, LIGHT_GRAY, header, 9, DARK_BLUE, True, BRIGHT_BLUE)
        x += width + 0.1

    # Table rows
    rows = [
        ("Layer 1", "YEARS, TEST_GROUP (3 items)", "Experiment Metadata table", "Table built", OCEAN),
        ("Layer 2", "CAMPAIGN_METADATA (18 items)", "Mnemonic Mapping v2", "MM v2 has metric fields", WARM_YELLOW),
        ("Layer 3", "SUCCESS_DEFINITIONS (24 items)", "GitHub or Semantic layer", "Library exists", WARM_YELLOW),
        ("Layer 4", "Source paths (8 items)", "Client Journey layer", "Layer built", OCEAN),
    ]

    for i, (layer, hardcoded, swaps_to, trigger, color) in enumerate(rows):
        y = 5.56 + (i * 0.32)
        add_box(slide, 0.35, y, 1.2, 0.28, color, layer, 8, WHITE if color == OCEAN else DARK_BLUE, True)
        add_text(slide, 1.65, y, 3.5, 0.28, hardcoded, 8, DARK_BLUE)
        add_text(slide, 5.25, y, 3.5, 0.28, swaps_to, 8, DARK_BLUE)
        add_text(slide, 8.85, y, 3.8, 0.28, trigger, 8, DARK_BLUE)

    # Total
    add_box(slide, 0.2, 6.88, 12.9, 0.4, SUNBURST,
            "TOTAL: 53 hardcoded items → Dynamic queries when infrastructure ready | ENGINE CORE: Does NOT change",
            10, DARK_BLUE, True)

    return slide

def create_page4_next_steps(prs):
    """Page 4: Next Steps"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text(slide, 0.3, 0.1, 12.7, 0.35, "Vintage Automation - Next Steps",
             font_size=26, bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.3, 0.4, 12.7, 0.25, "Strategic roadmap and decisions to be made",
             font_size=12, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Three main boxes
    boxes = [
        ("1. Adding New Cohorts", OCEAN, "For existing campaigns\n\nTRACK A (Official):\n• Automated refresh with CIDM\n• New cohorts appear automatically\n\nTRACK B (In-House):\n• Re-run engine\n• Update HTML, deploy to SharePoint"),
        ("2. Expand Metrics", SUNBURST, "What else can we vintage?\n\nSUCCESS METRICS:\n• Primary / Secondary / Tertiary\n\nENGAGEMENT:\n• Open rate, Click rate curves\n\nSEGMENT BREAKDOWNS:\n• By channel, by segment\n\n[TO BE FIGURED OUT]"),
        ("3. Hosting & Technology", TUNDRA, "Where does this live?\n\nOption A: SharePoint\n• Quick, manual refresh\n\nOption B: Tableau/CIDM\n• Official, governed\n\nOption C: Snowflake/New\n• Modern, scalable\n\n[DECISION PENDING]"),
    ]

    for i, (title, color, content) in enumerate(boxes):
        x = 0.3 + (i * 4.35)
        add_box(slide, x, 0.75, 4.15, 2.7, WHITE, "", border_color=color)
        add_box(slide, x, 0.75, 4.15, 0.3, color, title, 11, WHITE if color != SUNBURST else DARK_BLUE, True)
        add_text(slide, x + 0.1, 1.1, 3.95, 2.3, content, 9, DARK_BLUE)

    # Decisions section
    add_box(slide, 0.3, 3.55, 12.7, 1.5, WHITE, "", border_color=DARK_BLUE)
    add_box(slide, 0.3, 3.55, 12.7, 0.28, DARK_BLUE, "DECISIONS TO MAKE", 11, WHITE, True)

    decisions = [
        ("Refresh Process", "• How often?\n• Who triggers?\n• Automated vs manual?"),
        ("Metric Priority", "• What secondary metrics?\n• Which engagement?\n• Director input needed"),
        ("Technology", "• SharePoint enough?\n• Move to Snowflake?\n• Org direction?"),
        ("Track A vs B", "• Both parallel?\n• When does A take over?\n• Does B sunset?"),
    ]

    for i, (title, content) in enumerate(decisions):
        x = 0.45 + (i * 3.15)
        add_box(slide, x, 3.88, 3.0, 0.25, BRIGHT_BLUE, title, 10, WHITE, True)
        add_text(slide, x + 0.1, 4.18, 2.9, 0.8, content, 9, DARK_BLUE)

    # Key message
    add_box(slide, 0.3, 5.15, 12.7, 0.65, DARK_BLUE,
            "KEY MESSAGE\n\nEngine is built. Now we need to decide: how do we refresh, what else do we measure, and where does it live?",
            11, WHITE, True)

    # Status section
    add_box(slide, 0.3, 5.9, 12.7, 1.0, WHITE, "", border_color=BRIGHT_BLUE)
    add_box(slide, 0.3, 5.9, 12.7, 0.25, BRIGHT_BLUE, "CURRENT STATUS", 10, WHITE, True)

    add_box(slide, 0.5, 6.2, 4.0, 0.6, TUNDRA, "✓ DONE\n• Engine built\n• 6 pilots configured\n• Channel detection", 9, WHITE)
    add_box(slide, 4.7, 6.2, 4.0, 0.6, SUNBURST, "◐ IN PROGRESS\n• Track B dashboard\n• Success Library (4)\n• Swap points documented", 9, DARK_BLUE)
    add_box(slide, 8.9, 6.2, 4.0, 0.6, GRAY, "○ TO DO\n• Track A alignment\n• Hosting decision\n• Metric expansion", 9, WHITE)

    return slide

def create_highlevel_slide(prs):
    """Create a high-level one-page summary"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text(slide, 0.3, 0.1, 9, 0.4, "VINTAGE AUTOMATION",
             font_size=32, bold=True, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.3, 0.45, 9, 0.25, "Implementing the SuperFact Framework",
             font_size=14, font_color=DARK_BLUE, align=PP_ALIGN.CENTER)

    # Left: SuperFact Layers
    add_box(slide, 0.2, 0.85, 3.0, 3.6, LIGHT_GRAY, "", border_color=DARK_BLUE)
    add_text(slide, 0.3, 0.9, 2.8, 0.25, "SUPERFACT LAYERS", 11, DARK_BLUE, True, PP_ALIGN.CENTER)

    layers = [
        ("Layer 1: Experiment", "FROM DATA", OCEAN, WHITE),
        ("Layer 2: Campaign", "SWAP POINT", WARM_YELLOW, DARK_BLUE),
        ("Layer 3: Success", "SWAP POINT", WARM_YELLOW, DARK_BLUE),
        ("Layer 4: Client", "FROM DATA", OCEAN, WHITE),
    ]

    for i, (name, tag, color, fg) in enumerate(layers):
        y = 1.2 + (i * 0.75)
        add_box(slide, 0.35, y, 2.0, 0.55, color, name, 9, fg, True)
        add_box(slide, 2.45, y + 0.15, 0.65, 0.25, color, tag, 7, fg)

    add_text(slide, 0.35, 4.2, 2.7, 0.2, "L1 & L4: From Data NOW\nL2 & L3: Swap when ready", 8, DARK_BLUE, False, PP_ALIGN.CENTER)

    # Center: Engine
    add_box(slide, 3.4, 0.85, 3.2, 3.0, DARK_BLUE)
    add_text(slide, 3.5, 1.0, 3.0, 0.4, "VINTAGE\nAUTOMATION\nENGINE", 14, WHITE, True, PP_ALIGN.CENTER)

    add_box(slide, 3.55, 1.7, 2.9, 1.0, WHITE, "", border_color=OCEAN)
    add_text(slide, 3.65, 1.75, 2.7, 0.2, "6 Pilot Campaigns", 10, DARK_BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, 3.65, 2.0, 2.7, 0.6, "VCN   VDA   VDT\nVUI   VUT   VAW", 12, DARK_BLUE, True, PP_ALIGN.CENTER)

    add_box(slide, 3.55, 2.85, 2.9, 0.85, OCEAN)
    add_text(slide, 3.65, 2.9, 2.7, 0.75, "Produces:\nVintage Curves\nLift + CI", 10, WHITE, True, PP_ALIGN.CENTER)

    # Arrows to tracks
    add_text(slide, 6.7, 1.5, 0.5, 0.3, "→", 20, BRIGHT_BLUE, True)
    add_text(slide, 6.7, 2.8, 0.5, 0.3, "→", 20, TUNDRA, True)

    # Right: Tracks
    add_box(slide, 7.0, 0.85, 2.3, 1.4, BRIGHT_BLUE)
    add_text(slide, 7.1, 0.95, 2.1, 0.3, "TRACK A\nOfficial", 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 7.1, 1.4, 2.1, 0.7, "• Tableau/CIDM\n• Governed\n• Official", 9, WHITE, False, PP_ALIGN.LEFT)

    add_box(slide, 7.0, 2.4, 2.3, 1.4, TUNDRA)
    add_text(slide, 7.1, 2.5, 2.1, 0.3, "TRACK B\nIn-House", 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 7.1, 2.95, 2.1, 0.7, "• HTML/Plotly\n• SharePoint\n• Ready NOW", 9, WHITE, False, PP_ALIGN.LEFT)

    # Virtuous Cycle Banner
    add_box(slide, 0.2, 4.55, 9.1, 1.1, SUNBURST, "", border_color=DARK_BLUE)
    add_text(slide, 0.3, 4.6, 8.9, 0.25, "VIRTUOUS CYCLE: Success Library Growth", 12, DARK_BLUE, True, PP_ALIGN.CENTER)

    steps = ["New\nCampaign", "Define\nSuccess", "Add to\nLibrary", "Library\nGrows", "More\nGovernance", "Next is\nFASTER"]
    for i, step in enumerate(steps):
        x = 0.4 + (i * 1.5)
        add_box(slide, x, 4.9, 1.1, 0.6, DARK_BLUE, step, 8, WHITE, True)
        if i < 5:
            add_text(slide, x + 1.15, 5.1, 0.3, 0.2, "→", 12, DARK_BLUE, True)

    # Next Steps sidebar
    add_box(slide, 9.5, 0.85, 3.6, 5.8, WHITE, "", border_color=DARK_BLUE)
    add_box(slide, 9.5, 0.85, 3.6, 0.28, DARK_BLUE, "NEXT STEPS", 11, WHITE, True)

    add_box(slide, 9.6, 1.2, 3.4, 0.22, OCEAN, "1. Adding New Cohorts", 9, WHITE, True)
    add_text(slide, 9.7, 1.45, 3.3, 0.4, "• Track A: Automated with CIDM\n• Track B: Re-run, update HTML", 8, DARK_BLUE)

    add_box(slide, 9.6, 1.95, 3.4, 0.22, SUNBURST, "2. Expand Metrics", 9, DARK_BLUE, True)
    add_text(slide, 9.7, 2.2, 3.3, 0.5, "• Primary/Secondary/Tertiary\n• Email engagement curves\n• Segment breakdowns [TBD]", 8, DARK_BLUE)

    add_box(slide, 9.6, 2.8, 3.4, 0.22, TUNDRA, "3. Hosting Decision", 9, WHITE, True)
    add_text(slide, 9.7, 3.05, 3.3, 0.5, "• SharePoint (quick)\n• Tableau/CIDM (official)\n• Snowflake (modern) [TBD]", 8, DARK_BLUE)

    add_box(slide, 9.6, 3.65, 3.4, 0.22, DARK_BLUE, "DECISIONS PENDING", 9, WHITE, True)
    add_text(slide, 9.7, 3.9, 3.3, 0.5, "• Refresh cadence\n• Metric prioritization\n• Track A vs B balance", 8, DARK_BLUE)

    # Status
    add_box(slide, 9.6, 4.5, 3.4, 0.7, LIGHT_GRAY, "", border_color=BRIGHT_BLUE)
    add_box(slide, 9.7, 4.55, 1.0, 0.2, TUNDRA, "✓ Done", 7, WHITE)
    add_box(slide, 10.8, 4.55, 1.1, 0.2, SUNBURST, "◐ Progress", 7, DARK_BLUE)
    add_box(slide, 12.0, 4.55, 0.9, 0.2, GRAY, "○ To Do", 7, WHITE)
    add_text(slide, 9.7, 4.8, 3.2, 0.35, "Engine built → Dashboard ready → Decisions pending", 8, DARK_BLUE, False, PP_ALIGN.CENTER)

    # Legend
    add_box(slide, 0.2, 5.8, 6, 0.5, LIGHT_GRAY, "", border_color=DARK_BLUE)
    add_text(slide, 0.3, 5.85, 0.8, 0.2, "LEGEND:", 9, DARK_BLUE, True)
    add_box(slide, 1.1, 5.9, 0.2, 0.15, OCEAN)
    add_text(slide, 1.35, 5.85, 1.2, 0.2, "FROM DATA", 8, DARK_BLUE)
    add_box(slide, 2.6, 5.9, 0.2, 0.15, WARM_YELLOW)
    add_text(slide, 2.85, 5.85, 1.2, 0.2, "SWAP POINT", 8, DARK_BLUE)
    add_box(slide, 4.1, 5.9, 0.2, 0.15, SUNBURST)
    add_text(slide, 4.35, 5.85, 1.5, 0.2, "VIRTUOUS CYCLE", 8, DARK_BLUE)

    # Key message
    add_box(slide, 6.3, 5.8, 6.8, 0.5, DARK_BLUE,
            "Engine is built. Next: refresh process, expand metrics, decide on hosting.", 10, WHITE, True)

    return slide

def main():
    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    add_title_slide(prs, "VINTAGE AUTOMATION", "Implementing the SuperFact Framework")

    # Page 1: Big Picture
    create_page1_big_picture(prs)

    # Page 2: Virtuous Cycle
    create_page2_virtuous_cycle(prs)

    # Page 3: Swap Points
    create_page3_swap_points(prs)

    # Page 4: Next Steps
    create_page4_next_steps(prs)

    # High-level one-pager
    create_highlevel_slide(prs)

    # Save
    output_path = "/mnt/c/Users/andre/New_projects/Vintage/Vvd/pres/VINTAGE_AUTOMATION.pptx"
    prs.save(output_path)
    print(f"PowerPoint saved to: {output_path}")
    print("Slides created:")
    print("  1. Title Slide")
    print("  2. Big Picture - SuperFact Implementation")
    print("  3. Virtuous Cycle - Success Library Growth")
    print("  4. Swap Points - What Swaps To What")
    print("  5. Next Steps - Strategic Roadmap")
    print("  6. High-Level One-Pager (Director Summary)")

if __name__ == "__main__":
    main()
