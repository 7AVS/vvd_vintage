"""
DrawIO to PowerPoint Converter
Converts a layout specification JSON to .pptx with native editable shapes.

Usage:
    python converter.py layout.json [output.pptx]

Inputs:
    - layout.json: Output from the Layout Agent (element positions and styles)
    - config.json: Configuration with slide dimensions and offset (optional)

Output:
    - PowerPoint file with editable shapes
"""

import sys
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE


class Converter:
    def __init__(self, layout_path: str, config_path: str = None):
        self.layout_path = Path(layout_path)

        # Load layout spec
        with open(layout_path, 'r') as f:
            self.layout = json.load(f)

        # Load config (optional)
        self.config = {}
        if config_path and Path(config_path).exists():
            with open(config_path, 'r') as f:
                self.config = json.load(f)

        # Get dimensions from config or layout metadata
        slide_config = self.config.get('slide', {})
        layout_meta = self.layout.get('metadata', {})
        target = layout_meta.get('target_canvas', {})

        self.slide_width = slide_config.get('width', target.get('width', 1200))
        self.slide_height = slide_config.get('height', target.get('height', 675))
        self.offset_x = slide_config.get('offset_x', 0)
        self.offset_y = slide_config.get('offset_y', 0)

        # Scale: layout units to inches (assuming 16:9 slide = 13.333" x 7.5")
        self.scale = Inches(13.333) / self.slide_width

        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def convert(self, output_path: str = None):
        """Main conversion method."""
        if output_path is None:
            output_path = self.layout_path.with_suffix('.pptx')

        # Add blank slide
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # Process each element
        for elem in self.layout.get('elements', []):
            self._add_element(slide, elem)

        self.prs.save(output_path)
        print(f"Saved: {output_path}")
        return output_path

    def _scale(self, value: float) -> int:
        """Scale layout units to EMUs."""
        return int(value * self.scale)

    def _add_element(self, slide, elem: dict):
        """Add an element to the slide based on its type."""
        elem_type = elem.get('type', 'box')

        if elem_type == 'text':
            self._add_textbox(slide, elem)
        elif elem_type in ('box', 'container', 'badge'):
            self._add_shape(slide, elem)
        elif elem_type == 'line':
            self._add_line(slide, elem)
        elif elem_type == 'arrow':
            self._add_arrow(slide, elem)

    def _add_textbox(self, slide, elem: dict):
        """Add a text box."""
        left = self._scale(elem.get('x', 0) + self.offset_x)
        top = self._scale(elem.get('y', 0) + self.offset_y)
        width = self._scale(elem.get('width', 100))
        height = self._scale(elem.get('height', 30))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        value = elem.get('value', '')
        if value:
            p = tf.paragraphs[0]
            p.text = value
            self._apply_text_style(p, elem)

        tf.anchor = MSO_ANCHOR.MIDDLE

    def _add_shape(self, slide, elem: dict):
        """Add a shape (rectangle, rounded rectangle)."""
        left = self._scale(elem.get('x', 0) + self.offset_x)
        top = self._scale(elem.get('y', 0) + self.offset_y)
        width = self._scale(elem.get('width', 100))
        height = self._scale(elem.get('height', 50))

        # Determine shape type
        if elem.get('rounded', False):
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        else:
            shape_type = MSO_SHAPE.RECTANGLE

        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        # Apply fill
        fill_color = elem.get('fillColor', '#FFFFFF')
        if fill_color and fill_color.lower() not in ('none', 'transparent', ''):
            shape.fill.solid()
            rgb = self._hex_to_rgb(fill_color)
            if rgb:
                shape.fill.fore_color.rgb = rgb
        else:
            shape.fill.background()

        # Apply border
        stroke_color = elem.get('strokeColor')
        stroke_width = elem.get('strokeWidth', 1)

        if stroke_color and stroke_color.lower() not in ('none', 'transparent', ''):
            rgb = self._hex_to_rgb(stroke_color)
            if rgb:
                shape.line.color.rgb = rgb
            shape.line.width = Pt(float(stroke_width))
        else:
            shape.line.fill.background()

        # Dashed border
        if elem.get('dashed', False):
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # Add text if present
        value = elem.get('value', '')
        if value:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            p = tf.paragraphs[0]
            p.text = value
            self._apply_text_style(p, elem)
            tf.anchor = MSO_ANCHOR.MIDDLE

    def _add_line(self, slide, elem: dict):
        """Add a horizontal line."""
        left = self._scale(elem.get('x', 0) + self.offset_x)
        top = self._scale(elem.get('y', 0) + self.offset_y)
        width = self._scale(elem.get('width', 100))

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))

        stroke_color = elem.get('strokeColor', '#CCCCCC')
        rgb = self._hex_to_rgb(stroke_color)
        if rgb:
            line.fill.solid()
            line.fill.fore_color.rgb = rgb
        line.line.fill.background()

    def _add_arrow(self, slide, elem: dict):
        """Add an arrow/connector."""
        x1 = self._scale(elem.get('x1', 0) + self.offset_x)
        y1 = self._scale(elem.get('y1', 0) + self.offset_y)
        x2 = self._scale(elem.get('x2', 0) + self.offset_x)
        y2 = self._scale(elem.get('y2', 0) + self.offset_y)

        left = min(x1, x2)
        top = min(y1, y2)
        width = abs(x2 - x1) or Emu(1)
        height = abs(y2 - y1) or Emu(1)

        connector = slide.shapes.add_connector(
            1, left, top, left + width, top + height
        )

        stroke_color = elem.get('strokeColor', '#000000')
        stroke_width = elem.get('strokeWidth', 1)

        rgb = self._hex_to_rgb(stroke_color)
        if rgb:
            connector.line.color.rgb = rgb
        connector.line.width = Pt(float(stroke_width))

    def _apply_text_style(self, paragraph, elem: dict):
        """Apply text styling to a paragraph."""
        font = paragraph.font

        # Font size
        font_size = elem.get('fontSize', 11)
        font.size = Pt(float(font_size))

        # Font color
        font_color = elem.get('fontColor', '#000000')
        if font_color:
            rgb = self._hex_to_rgb(font_color)
            if rgb:
                font.color.rgb = rgb

        # Font style (1=bold, 2=italic, 3=bold+italic)
        font_style = elem.get('fontStyle', 0)
        if isinstance(font_style, str):
            font_style = int(font_style) if font_style.isdigit() else 0
        font.bold = font_style in (1, 3)
        font.italic = font_style in (2, 3)

        # Font family
        font_family = elem.get('fontFamily')
        if font_family:
            font.name = font_family

        # NO EFFECTS - clean text only
        font.shadow = None
        font.underline = False

        # Text alignment
        align = elem.get('align', 'center')
        if align == 'left':
            paragraph.alignment = PP_ALIGN.LEFT
        elif align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.CENTER

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        if not hex_color or hex_color.lower() in ('none', 'transparent', ''):
            return None
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        return RGBColor(0, 0, 0)


def main():
    if len(sys.argv) < 2:
        print("Usage: python converter.py layout.json [output.pptx]")
        print("\nConverts a layout specification JSON to PowerPoint.")
        sys.exit(1)

    layout_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    # Look for config.json in same directory
    config_file = Path(layout_file).parent / 'config.json'

    converter = Converter(layout_file, str(config_file) if config_file.exists() else None)
    converter.convert(output_file)


if __name__ == '__main__':
    main()
