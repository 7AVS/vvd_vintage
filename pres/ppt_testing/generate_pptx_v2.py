"""
Vintage Automation PowerPoint Generator v2
Minimalist approach - strategic color usage
Based on VISUAL_FRAMEWORK.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# === COLOR PALETTE (Minimalist) ===
# Primary
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BLUE = RGBColor(0x00, 0x31, 0x68)
MEDIUM_GRAY = RGBColor(0x9E, 0xA2, 0xA2)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# Accents (use sparingly)
OCEAN = RGBColor(0x00, 0x91, 0xDA)       # From Data
WARM_YELLOW = RGBColor(0xFF, 0xC7, 0x2C) # Swap Point
SUNBURST = RGBColor(0xFC, 0xA3, 0x11)    # Key emphasis
TUNDRA = RGBColor(0x07, 0xAF, 0xBF)      # Future/Growth


def set_shape_border(shape, color, width=1):
    """Set border color and width, no fill"""
    shape.fill.background()  # Transparent fill
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_shape_fill_and_border(shape, fill_color, border_color=None, border_width=1):
    """Set fill and optional border"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text,
                 font_size=11, font_color=DARK_BLUE, bold=False,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    """Add a simple text box"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color

    return textbox


def add_box_border_only(slide, left, top, width, height, border_color=BORDER_GRAY,
                        text="", font_size=10, font_color=DARK_BLUE, bold=False, border_width=1.5):
    """Add a box with border only (no fill) - primary style"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape_border(shape, border_color, border_width)

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color

    return shape


def add_box_filled(slide, left, top, width, height, fill_color,
                   text="", font_size=10, font_color=WHITE, bold=False, border_color=None):
    """Add a filled box - use sparingly for emphasis"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape_fill_and_border(shape, fill_color, border_color)

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color

    return shape


def add_tag(slide, left, top, text, color, font_color=None):
    """Add a small colored tag"""
    if font_color is None:
        # Auto-determine font color based on background
        font_color = WHITE if color in [OCEAN, DARK_BLUE, TUNDRA] else DARK_BLUE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(0.8), Inches(0.18))
    set_shape_fill_and_border(shape, color)

    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = font_color

    return shape


def add_arrow_text(slide, left, top, direction="down", color=MEDIUM_GRAY):
    """Add arrow as text (simple approach)"""
    arrows = {"down": "▼", "right": "→", "left": "←", "up": "▲"}
    arrow = arrows.get(direction, "→")

    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(0.3), Inches(0.3))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = arrow
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    return textbox


def add_line(slide, start_x, start_y, end_x, end_y, color=BORDER_GRAY, width=1):
    """Add a simple line"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start_x), Inches(start_y),
        Inches(end_x), Inches(end_y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


# === SLIDE CREATORS ===

def create_title_slide(prs):
    """Slide 1: Clean minimal title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title - large, centered
    add_text_box(slide, 0, 2.8, 13.33, 0.8, "VINTAGE AUTOMATION",
                 font_size=48, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, 0, 3.7, 13.33, 0.5, "Implementing the SuperFact Framework",
                 font_size=20, align=PP_ALIGN.CENTER, font_color=MEDIUM_GRAY)

    # Simple line separator
    add_line(slide, 5, 4.3, 8.33, 4.3, BORDER_GRAY, 1)

    return slide


def create_big_picture(prs):
    """Slide 2: Big Picture - Flow diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text_box(slide, 0.5, 0.3, 12, 0.4, "The Big Picture",
                 font_size=28, bold=True, align=PP_ALIGN.LEFT)
    add_text_box(slide, 0.5, 0.7, 12, 0.3, "SuperFact layers feed the engine, engine produces outputs",
                 font_size=12, font_color=MEDIUM_GRAY)

    # === SUPERFACT LAYERS (top) ===
    add_text_box(slide, 0.5, 1.2, 3, 0.25, "SUPERFACT LAYERS", font_size=10, bold=True, font_color=MEDIUM_GRAY)

    layers = [
        ("Layer 1", "Experiment\nMetadata", OCEAN, "FROM DATA"),
        ("Layer 2", "Campaign\nMetadata", WARM_YELLOW, "SWAP POINT"),
        ("Layer 3", "Success\nLibrary", WARM_YELLOW, "SWAP POINT"),
        ("Layer 4", "Client\nJourney", OCEAN, "FROM DATA"),
    ]

    layer_width = 2.3
    layer_start = 0.5
    for i, (title, subtitle, color, tag) in enumerate(layers):
        x = layer_start + (i * 2.6)
        # Box with colored border only
        add_box_border_only(slide, x, 1.5, layer_width, 1.0, color, f"{title}\n{subtitle}", 10, DARK_BLUE)
        # Small tag below
        add_tag(slide, x + 0.75, 2.55, tag, color)

    # Arrow down to engine
    add_arrow_text(slide, 5.2, 2.85, "down", MEDIUM_GRAY)

    # === ENGINE (center) ===
    # Light gray background box
    add_box_filled(slide, 3.5, 3.2, 6.33, 1.8, LIGHT_GRAY, border_color=DARK_BLUE)
    add_text_box(slide, 3.7, 3.3, 5.9, 0.35, "VINTAGE AUTOMATION ENGINE",
                 font_size=14, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 3.7, 3.7, 5.9, 0.25, "vintage_all_in_one.py",
                 font_size=9, font_color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # Campaigns inside engine
    add_box_border_only(slide, 4.0, 4.05, 5.33, 0.8, BORDER_GRAY)
    add_text_box(slide, 4.1, 4.1, 5.1, 0.2, "6 Pilot Campaigns", font_size=9, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 4.1, 4.35, 5.1, 0.4, "VCN    VDA    VDT    VUI    VUT    VAW",
                 font_size=11, align=PP_ALIGN.CENTER)

    # Arrow down to outputs
    add_arrow_text(slide, 6.5, 5.1, "down", MEDIUM_GRAY)

    # === OUTPUT (bottom) ===
    add_text_box(slide, 5.5, 5.45, 2.5, 0.25, "CSV / HDFS", font_size=10, bold=True, align=PP_ALIGN.CENTER)

    # Two tracks
    add_arrow_text(slide, 4.2, 5.7, "down", MEDIUM_GRAY)
    add_arrow_text(slide, 8.8, 5.7, "down", MEDIUM_GRAY)

    # Track A - border only
    add_box_border_only(slide, 2.5, 6.0, 3.2, 1.2, BORDER_GRAY)
    add_text_box(slide, 2.6, 6.05, 3.0, 0.25, "TRACK A", font_size=11, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 2.6, 6.3, 3.0, 0.2, "Official", font_size=9, align=PP_ALIGN.CENTER, font_color=MEDIUM_GRAY)
    add_text_box(slide, 2.6, 6.55, 3.0, 0.6, "• Tableau / CIDM\n• Governed\n• Pending alignment",
                 font_size=9, align=PP_ALIGN.LEFT)

    # Track B - border only
    add_box_border_only(slide, 7.6, 6.0, 3.2, 1.2, BORDER_GRAY)
    add_text_box(slide, 7.7, 6.05, 3.0, 0.25, "TRACK B", font_size=11, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 7.7, 6.3, 3.0, 0.2, "In-House", font_size=9, align=PP_ALIGN.CENTER, font_color=MEDIUM_GRAY)
    add_text_box(slide, 7.7, 6.55, 3.0, 0.6, "• HTML / Plotly\n• SharePoint\n• Ready now",
                 font_size=9, align=PP_ALIGN.LEFT)

    # Legend - minimal
    add_text_box(slide, 10.5, 1.2, 2.5, 0.2, "LEGEND", font_size=8, bold=True, font_color=MEDIUM_GRAY)
    add_line(slide, 10.5, 1.45, 10.7, 1.45, OCEAN, 3)
    add_text_box(slide, 10.8, 1.35, 1.5, 0.2, "From Data", font_size=8)
    add_line(slide, 10.5, 1.7, 10.7, 1.7, WARM_YELLOW, 3)
    add_text_box(slide, 10.8, 1.6, 1.5, 0.2, "Swap Point", font_size=8)

    return slide


def create_virtuous_cycle(prs):
    """Slide 3: Virtuous Cycle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text_box(slide, 0.5, 0.3, 12, 0.4, "The Virtuous Cycle",
                 font_size=28, bold=True)
    add_text_box(slide, 0.5, 0.7, 12, 0.3,
                 "As we add campaigns, we build the Success Library. The library grows. Next campaign is faster.",
                 font_size=12, font_color=MEDIUM_GRAY)

    # === THE CYCLE (left side) ===
    add_box_border_only(slide, 0.5, 1.2, 5.5, 4.0, BORDER_GRAY)
    add_text_box(slide, 0.6, 1.3, 5.3, 0.25, "THE CYCLE", font_size=11, bold=True, font_color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    steps = [
        ("1. New Campaign Onboarded", False),
        ("2. Understand Success Metric", False),
        ("3. Document in Success Library", True),  # Key step - emphasized
        ("4. Success Library Grows", False),
        ("5. Governance Improves", False),
        ("6. Next Campaign is Faster", False),
    ]

    y = 1.65
    for i, (text, is_key) in enumerate(steps):
        if is_key:
            # Key step - sunburst border
            add_box_border_only(slide, 0.8, y, 4.9, 0.4, SUNBURST, text, 10, DARK_BLUE, True, 2)
        else:
            add_box_border_only(slide, 0.8, y, 4.9, 0.4, BORDER_GRAY, text, 10, DARK_BLUE)

        if i < 5:
            add_arrow_text(slide, 3.0, y + 0.4, "down", MEDIUM_GRAY)
        y += 0.55

    # Loop arrow indication
    add_text_box(slide, 5.0, 3.5, 0.4, 0.3, "↺", font_size=20, font_color=MEDIUM_GRAY)

    # === LIBRARY GROWTH (right side) ===
    add_box_border_only(slide, 6.2, 1.2, 6.8, 4.0, BORDER_GRAY)
    add_text_box(slide, 6.3, 1.3, 6.6, 0.25, "SUCCESS LIBRARY GROWTH", font_size=11, bold=True, font_color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    campaigns = [
        ("VCN", "Card Acquisition", ["card_acq"], [True]),
        ("VDA", "Card Acquisition", ["card_acq"], [False]),
        ("VDT", "Card Activation", ["card_acq", "card_act"], [False, True]),
        ("VUI", "Card Usage", ["card_acq", "card_act", "card_use"], [False, False, True]),
        ("VUT", "Wallet Provision", ["card_acq", "card_act", "card_use", "wallet"], [False, False, False, True]),
        ("VAW", "Wallet Provision", ["card_acq", "card_act", "card_use", "wallet"], [False, False, False, False]),
    ]

    y = 1.65
    for mne, need, metrics, is_new in campaigns:
        # Campaign label
        add_text_box(slide, 6.4, y, 2.5, 0.35, f"{mne}: {need}", font_size=9, bold=True)

        # Metrics as small tags
        x = 8.9
        for j, (metric, new) in enumerate(zip(metrics, is_new)):
            if new:
                add_tag(slide, x, y + 0.08, metric, TUNDRA)
            else:
                add_tag(slide, x, y + 0.08, metric, LIGHT_GRAY, MEDIUM_GRAY)
            x += 0.85

        y += 0.4

    # Result callout - the ONE emphasized element
    add_box_filled(slide, 6.2, 4.35, 6.8, 0.7, SUNBURST,
                   "RESULT: 6 campaigns, only 4 unique metrics\nFuture campaigns reuse existing metrics",
                   10, DARK_BLUE, True)

    # Key message at bottom
    add_text_box(slide, 0.5, 5.4, 12.3, 0.5,
                 "\"Every campaign onboarded enriches our metadata ecosystem. We're building the source of truth.\"",
                 font_size=12, font_color=DARK_BLUE, bold=False, align=PP_ALIGN.CENTER)

    return slide


def create_swap_points(prs):
    """Slide 4: Swap Points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text_box(slide, 0.5, 0.3, 12, 0.4, "Swap Points",
                 font_size=28, bold=True)
    add_text_box(slide, 0.5, 0.7, 12, 0.3,
                 "Hardcoded today, dynamic when infrastructure is ready. Engine core doesn't change.",
                 font_size=12, font_color=MEDIUM_GRAY)

    # === LAYER 2 ===
    add_box_border_only(slide, 0.5, 1.15, 12.3, 1.7, WARM_YELLOW, border_width=2)
    add_text_box(slide, 0.7, 1.25, 11.9, 0.25, "LAYER 2: Campaign Metadata", font_size=12, bold=True)
    add_text_box(slide, 0.7, 1.5, 11.9, 0.2, "\"What metric should we measure for this campaign?\"", font_size=9, font_color=MEDIUM_GRAY)

    # Today
    add_text_box(slide, 0.7, 1.8, 2, 0.2, "TODAY", font_size=9, bold=True, font_color=MEDIUM_GRAY)
    add_box_border_only(slide, 0.7, 2.0, 4.8, 0.75, BORDER_GRAY,
                        "CAMPAIGN_METADATA = {\n  \"VCN\": {\"primary_metric\": \"card_acquisition\"},\n  ...6 campaigns hardcoded\n}",
                        8, DARK_BLUE)

    # Arrow
    add_text_box(slide, 5.7, 2.2, 1, 0.3, "→", font_size=24, font_color=SUNBURST, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 5.6, 2.5, 1.2, 0.2, "SWAP", font_size=8, bold=True, font_color=SUNBURST, align=PP_ALIGN.CENTER)

    # Future
    add_text_box(slide, 6.8, 1.8, 2, 0.2, "FUTURE", font_size=9, bold=True, font_color=MEDIUM_GRAY)
    add_box_border_only(slide, 6.8, 2.0, 5.8, 0.75, TUNDRA,
                        "SELECT primary_metric, secondary_metric\nFROM CIDM_MNEMONIC_ATTRS\nWHERE mne = 'VCN'",
                        8, DARK_BLUE)

    # === LAYER 3 ===
    add_box_border_only(slide, 0.5, 3.0, 12.3, 1.7, WARM_YELLOW, border_width=2)
    add_text_box(slide, 0.7, 3.1, 11.9, 0.25, "LAYER 3: Success Definitions", font_size=12, bold=True)
    add_text_box(slide, 0.7, 3.35, 11.9, 0.2, "\"HOW do we calculate this metric?\"", font_size=9, font_color=MEDIUM_GRAY)

    # Today
    add_text_box(slide, 0.7, 3.6, 2, 0.2, "TODAY", font_size=9, bold=True, font_color=MEDIUM_GRAY)
    add_box_border_only(slide, 0.7, 3.8, 4.8, 0.8, BORDER_GRAY,
                        "SUCCESS_DEFINITIONS = {\n  \"card_acquisition\": {\n    \"table\": \"VISA_DR_CRD\", \"filters\": {...}\n  }... 4 metrics × 6 fields = 24 items\n}",
                        8, DARK_BLUE)

    # Arrow
    add_text_box(slide, 5.7, 4.0, 1, 0.3, "→", font_size=24, font_color=SUNBURST, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 5.6, 4.3, 1.2, 0.2, "SWAP", font_size=8, bold=True, font_color=SUNBURST, align=PP_ALIGN.CENTER)

    # Future
    add_text_box(slide, 6.8, 3.6, 2, 0.2, "FUTURE", font_size=9, bold=True, font_color=MEDIUM_GRAY)
    add_box_border_only(slide, 6.8, 3.8, 2.7, 0.8, TUNDRA,
                        "Option A:\n%run /success_library/\n      card_acquisition.py",
                        8, DARK_BLUE)
    add_box_border_only(slide, 9.7, 3.8, 2.9, 0.8, TUNDRA,
                        "Option B:\nSELECT * FROM\nsemantic.success_card_acq",
                        8, DARK_BLUE)

    # === SUMMARY TABLE ===
    add_text_box(slide, 0.5, 4.9, 12, 0.25, "SUMMARY", font_size=11, bold=True, font_color=MEDIUM_GRAY)

    # Table header
    headers = [("Layer", 1.2), ("Hardcoded", 3.5), ("Swaps To", 3.5), ("Trigger", 3.5)]
    x = 0.5
    for header, width in headers:
        add_box_filled(slide, x, 5.15, width, 0.3, LIGHT_GRAY, header, 9, DARK_BLUE, True, BORDER_GRAY)
        x += width + 0.1

    # Table rows
    rows = [
        ("L1", "Years, test groups (3)", "Experiment Metadata table", "Table built", OCEAN),
        ("L2", "Campaign config (18)", "Mnemonic Mapping v2", "MM v2 has metric fields", WARM_YELLOW),
        ("L3", "Success definitions (24)", "GitHub or Semantic layer", "Library exists", WARM_YELLOW),
        ("L4", "Source paths (8)", "Client Journey layer", "Layer built", OCEAN),
    ]

    y = 5.5
    for layer, hardcoded, swaps_to, trigger, color in rows:
        add_box_border_only(slide, 0.5, y, 1.2, 0.3, color, layer, 9, DARK_BLUE, True, 2)
        add_text_box(slide, 1.8, y + 0.05, 3.4, 0.25, hardcoded, 8)
        add_text_box(slide, 5.4, y + 0.05, 3.4, 0.25, swaps_to, 8)
        add_text_box(slide, 9.0, y + 0.05, 3.4, 0.25, trigger, 8)
        y += 0.35

    # Total - emphasized
    add_box_border_only(slide, 0.5, 6.95, 12.3, 0.35, SUNBURST,
                        "TOTAL: 53 hardcoded items → Dynamic queries when infrastructure ready",
                        10, DARK_BLUE, True, 2)

    return slide


def create_next_steps(prs):
    """Slide 5: Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text_box(slide, 0.5, 0.3, 12, 0.4, "Next Steps",
                 font_size=28, bold=True)
    add_text_box(slide, 0.5, 0.7, 12, 0.3,
                 "Engine is built. Now: how do we refresh, what else do we measure, where does it live?",
                 font_size=12, font_color=MEDIUM_GRAY)

    # === THREE PILLARS ===
    pillars = [
        ("1. Adding New Cohorts",
         "For existing campaigns\n\nTrack A:\n• Automated refresh with CIDM\n• New cohorts appear automatically\n\nTrack B:\n• Re-run engine\n• Update HTML, deploy to SharePoint",
         OCEAN),
        ("2. Expand Metrics",
         "What else can we vintage?\n\n• Primary / Secondary / Tertiary\n• Email engagement curves\n• Segment breakdowns\n• Other comparisons\n\n[TO BE FIGURED OUT]",
         SUNBURST),
        ("3. Hosting Decision",
         "Where does this live?\n\n• SharePoint (quick, manual)\n• Tableau/CIDM (official, governed)\n• Snowflake/New (modern, scalable)\n\n[DECISION PENDING]",
         TUNDRA),
    ]

    x = 0.5
    for title, content, accent_color in pillars:
        # Pillar box
        add_box_border_only(slide, x, 1.1, 4.0, 3.0, BORDER_GRAY)
        # Header with accent
        add_box_filled(slide, x, 1.1, 4.0, 0.4, LIGHT_GRAY, title, 11, DARK_BLUE, True, accent_color)
        # Content
        add_text_box(slide, x + 0.15, 1.55, 3.7, 2.5, content, 9)
        x += 4.2

    # === DECISIONS ===
    add_line(slide, 0.5, 4.3, 12.8, 4.3, BORDER_GRAY, 1)
    add_text_box(slide, 0.5, 4.4, 12, 0.25, "DECISIONS PENDING", font_size=11, bold=True, font_color=MEDIUM_GRAY)

    decisions = [
        "Refresh cadence - how often?",
        "Metric prioritization - what matters?",
        "Technology platform - SharePoint vs Snowflake?",
        "Track A vs B - both parallel or one takes over?",
    ]

    add_text_box(slide, 0.5, 4.7, 12.3, 1.0,
                 "  •  ".join(decisions),
                 font_size=10, align=PP_ALIGN.CENTER)

    # === STATUS ===
    add_line(slide, 0.5, 5.3, 12.8, 5.3, BORDER_GRAY, 1)
    add_text_box(slide, 0.5, 5.4, 2, 0.25, "STATUS", font_size=11, bold=True, font_color=MEDIUM_GRAY)

    # Status boxes
    add_box_filled(slide, 2.5, 5.4, 3.0, 0.5, LIGHT_GRAY, "✓ DONE\nEngine built, 6 pilots, channel detection", 8, DARK_BLUE, border_color=TUNDRA)
    add_box_filled(slide, 5.7, 5.4, 3.0, 0.5, LIGHT_GRAY, "◐ IN PROGRESS\nDashboard ready, Success Library", 8, DARK_BLUE, border_color=SUNBURST)
    add_box_filled(slide, 8.9, 5.4, 3.5, 0.5, LIGHT_GRAY, "○ TO DO\nTrack A alignment, hosting, metrics", 8, DARK_BLUE, border_color=MEDIUM_GRAY)

    return slide


def create_one_pager(prs):
    """Slide 6: High-level one-pager summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text_box(slide, 0.3, 0.15, 9, 0.35, "VINTAGE AUTOMATION",
                 font_size=24, bold=True)
    add_text_box(slide, 0.3, 0.45, 9, 0.2, "Implementing the SuperFact Framework",
                 font_size=11, font_color=MEDIUM_GRAY)

    # === MAIN FLOW (left to right) ===
    # Layers
    add_box_border_only(slide, 0.3, 0.85, 2.2, 2.0, BORDER_GRAY)
    add_text_box(slide, 0.4, 0.9, 2.0, 0.2, "LAYERS", font_size=9, bold=True, font_color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    layers = [("L1", OCEAN), ("L2", WARM_YELLOW), ("L3", WARM_YELLOW), ("L4", OCEAN)]
    y = 1.15
    for layer, color in layers:
        add_box_border_only(slide, 0.5, y, 1.8, 0.35, color, layer, 9, DARK_BLUE)
        y += 0.42

    # Arrow
    add_text_box(slide, 2.6, 1.7, 0.4, 0.3, "→", font_size=18, font_color=MEDIUM_GRAY)

    # Engine
    add_box_filled(slide, 3.1, 0.85, 2.5, 2.0, LIGHT_GRAY, border_color=DARK_BLUE)
    add_text_box(slide, 3.2, 0.95, 2.3, 0.35, "ENGINE", font_size=10, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 3.2, 1.3, 2.3, 0.3, "6 Campaigns\nVintage Curves", font_size=9, align=PP_ALIGN.CENTER)
    add_text_box(slide, 3.2, 1.8, 2.3, 0.3, "Lift + CI", font_size=9, align=PP_ALIGN.CENTER, font_color=MEDIUM_GRAY)

    # Arrow
    add_text_box(slide, 5.7, 1.7, 0.4, 0.3, "→", font_size=18, font_color=MEDIUM_GRAY)

    # Outputs
    add_box_border_only(slide, 6.2, 0.85, 2.5, 0.9, BORDER_GRAY)
    add_text_box(slide, 6.3, 0.9, 2.3, 0.2, "TRACK A", font_size=9, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 6.3, 1.1, 2.3, 0.5, "Official\nTableau/CIDM", font_size=8, align=PP_ALIGN.CENTER, font_color=MEDIUM_GRAY)

    add_box_border_only(slide, 6.2, 1.85, 2.5, 0.9, BORDER_GRAY)
    add_text_box(slide, 6.3, 1.9, 2.3, 0.2, "TRACK B", font_size=9, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 6.3, 2.1, 2.3, 0.5, "In-House\nHTML/Plotly", font_size=8, align=PP_ALIGN.CENTER, font_color=MEDIUM_GRAY)

    # === VIRTUOUS CYCLE (banner) ===
    add_box_border_only(slide, 0.3, 3.0, 8.4, 0.7, SUNBURST, border_width=2)
    add_text_box(slide, 0.4, 3.05, 8.2, 0.2, "VIRTUOUS CYCLE", font_size=9, bold=True, font_color=SUNBURST, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.4, 3.3, 8.2, 0.35, "New Campaign → Define Success → Add to Library → Library Grows → Next is Faster",
                 font_size=9, align=PP_ALIGN.CENTER)

    # === SWAP POINTS ===
    add_text_box(slide, 0.3, 3.85, 8.4, 0.2, "SWAP POINTS", font_size=9, bold=True, font_color=MEDIUM_GRAY)
    add_text_box(slide, 0.3, 4.05, 8.4, 0.3, "L2 & L3 hardcoded today (53 items) → Dynamic queries when infrastructure ready",
                 font_size=9)

    # === NEXT STEPS (right sidebar) ===
    add_box_border_only(slide, 9.0, 0.85, 4.0, 3.5, BORDER_GRAY)
    add_text_box(slide, 9.1, 0.9, 3.8, 0.25, "NEXT STEPS", font_size=10, bold=True, align=PP_ALIGN.CENTER)

    steps = [
        ("1. Cohorts", "Add new cohorts for existing campaigns", OCEAN),
        ("2. Metrics", "Expand what we measure [TBD]", SUNBURST),
        ("3. Hosting", "Where does this live? [TBD]", TUNDRA),
    ]

    y = 1.2
    for title, desc, color in steps:
        add_box_border_only(slide, 9.15, y, 3.7, 0.25, color, title, 9, DARK_BLUE, True)
        add_text_box(slide, 9.15, y + 0.28, 3.7, 0.25, desc, 8, MEDIUM_GRAY)
        y += 0.6

    # Decisions
    add_text_box(slide, 9.15, 3.1, 3.7, 0.2, "DECISIONS PENDING:", font_size=8, bold=True, font_color=MEDIUM_GRAY)
    add_text_box(slide, 9.15, 3.3, 3.7, 0.9, "• Refresh cadence\n• Metric priority\n• Track A vs B balance",
                 font_size=8)

    # === LEGEND (bottom) ===
    add_line(slide, 0.3, 4.5, 12.7, 4.5, BORDER_GRAY, 0.5)
    add_text_box(slide, 0.3, 4.55, 1, 0.2, "LEGEND:", font_size=8, font_color=MEDIUM_GRAY)
    add_line(slide, 1.3, 4.65, 1.5, 4.65, OCEAN, 3)
    add_text_box(slide, 1.55, 4.55, 1.2, 0.2, "From Data", font_size=8)
    add_line(slide, 2.8, 4.65, 3.0, 4.65, WARM_YELLOW, 3)
    add_text_box(slide, 3.05, 4.55, 1.2, 0.2, "Swap Point", font_size=8)
    add_line(slide, 4.3, 4.65, 4.5, 4.65, SUNBURST, 3)
    add_text_box(slide, 4.55, 4.55, 1.5, 0.2, "Key Emphasis", font_size=8)

    return slide


def main():
    """Generate the PowerPoint"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create slides
    create_title_slide(prs)
    create_big_picture(prs)
    create_virtuous_cycle(prs)
    create_swap_points(prs)
    create_next_steps(prs)
    create_one_pager(prs)

    # Save
    output_path = "/mnt/c/Users/andre/New_projects/Vintage/Vvd/pres/ppt_testing/VINTAGE_AUTOMATION_v2.pptx"
    prs.save(output_path)

    print(f"✓ PowerPoint saved: {output_path}")
    print("\nSlides:")
    print("  1. Title")
    print("  2. Big Picture (flow diagram)")
    print("  3. Virtuous Cycle (circular + growth)")
    print("  4. Swap Points (today → future)")
    print("  5. Next Steps (three pillars)")
    print("  6. One-Pager Summary")
    print("\nStyle: Minimalist - white backgrounds, color used sparingly for emphasis")


if __name__ == "__main__":
    main()
